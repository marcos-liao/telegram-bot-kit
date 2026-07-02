#!/usr/bin/env python3
"""
Document text-extraction helpers used by the per-chat document upload/QA feature.

Supports PDF, DOCX, XLSX, PPTX, TXT/MD, and images (via OCR fallback for
scanned PDFs and image files). Pure extraction + chunking — no storage,
no embedding, no knowledge-base scanning.
"""
import os

DOC_OCR = os.environ.get("DOC_OCR", "1") not in ("0", "false", "no", "")
OCR_LANG = os.environ.get("OCR_LANG", "eng")
OCR_MAX_PAGES = int(os.environ.get("OCR_MAX_PAGES", "40"))   # cap pages OCR'd per PDF
PDF_MIN_TEXT = int(os.environ.get("PDF_MIN_TEXT", "120"))    # below this -> treat PDF as scanned

DOC_CHUNK_SIZE = int(os.environ.get("DOC_CHUNK_SIZE", "900"))
DOC_CHUNK_OVERLAP = int(os.environ.get("DOC_CHUNK_OVERLAP", "150"))

SUPPORTED = (".pdf", ".docx", ".xlsx", ".pptx", ".txt", ".md", ".png", ".jpg", ".jpeg")


# ----------------------------- extractors -----------------------------
def _ocr_image(img):
    import pytesseract

    return pytesseract.image_to_string(img, lang=OCR_LANG)


def _ocr_pdf(path):
    """Render PDF pages to images and OCR them (for scanned PDFs)."""
    import io
    import fitz  # PyMuPDF
    from PIL import Image

    parts = []
    doc = fitz.open(path)
    try:
        for i, page in enumerate(doc):
            if i >= OCR_MAX_PAGES:
                break
            pix = page.get_pixmap(dpi=200)
            with Image.open(io.BytesIO(pix.tobytes("png"))) as img:
                parts.append(_ocr_image(img))
    finally:
        doc.close()
    return "\n".join(parts)


def extract_pdf(path):
    from pypdf import PdfReader

    reader = PdfReader(path)
    text = "\n".join((page.extract_text() or "") for page in reader.pages)
    # Little/no text layer -> likely a scanned PDF -> OCR fallback.
    if DOC_OCR and len(text.strip()) < PDF_MIN_TEXT:
        ocr = _ocr_pdf(path)
        if len(ocr.strip()) > len(text.strip()):
            return ocr
    return text


def extract_image(path):
    if not DOC_OCR:
        return ""
    from PIL import Image

    with Image.open(path) as img:
        return _ocr_image(img)


def extract_docx(path):
    import docx

    d = docx.Document(path)
    parts = [p.text for p in d.paragraphs if p.text.strip()]
    for table in d.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def extract_xlsx(path):
    from openpyxl import load_workbook

    wb = load_workbook(path, read_only=True, data_only=True)
    parts = []
    for ws in wb.worksheets:
        parts.append(f"# Sheet: {ws.title}")
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None and str(c).strip()]
            if cells:
                parts.append(" | ".join(cells))
    wb.close()
    return "\n".join(parts)


def extract_pptx(path):
    from pptx import Presentation

    prs = Presentation(path)
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        parts.append(f"# Slide {i}")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    txt = "".join(run.text for run in para.runs).strip()
                    if txt:
                        parts.append(txt)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells if c.text.strip()]
                    if cells:
                        parts.append(" | ".join(cells))
    return "\n".join(parts)


def extract_txt(path):
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()


def extract(path):
    ext = os.path.splitext(path)[1].lower()
    if ext == ".pdf":
        return extract_pdf(path)
    if ext == ".docx":
        return extract_docx(path)
    if ext == ".xlsx":
        return extract_xlsx(path)
    if ext == ".pptx":
        return extract_pptx(path)
    if ext in (".txt", ".md"):
        return extract_txt(path)
    if ext in (".png", ".jpg", ".jpeg"):
        return extract_image(path)
    return ""


def chunk_text(text, size=DOC_CHUNK_SIZE, overlap=DOC_CHUNK_OVERLAP):
    text = " ".join(text.split())  # normalize whitespace
    if not text:
        return []
    chunks, i, n = [], 0, len(text)
    while i < n:
        chunks.append(text[i : i + size])
        i += max(1, size - overlap)
    return chunks
