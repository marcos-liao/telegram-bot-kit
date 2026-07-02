#!/usr/bin/env python3
"""
Telegram Bot Kit - AI assistant bot powered by any OpenAI-compatible LLM + SearXNG web search.

Features:
- Conversation memory per chat (SQLite): recent history + rolling summary.
- Automatic web search via tool-calling (model decides when to search).
- Runs 24x7 as a systemd service (long polling, no public webhook needed).

Config is read from environment (see .env.example):
  TELEGRAM_TOKEN   (required)  Telegram bot token from @BotFather
  LLM_BASE_URL     OpenAI-compatible base URL, default http://127.0.0.1:11434/v1
  LLM_MODEL        default llama3.1
  SEARXNG_URL      default http://127.0.0.1:8888
  DB_PATH          default ./memory.db
"""

import os
import re
import json
import time
import uuid
import base64
import asyncio
import signal
import ipaddress
import subprocess
import logging
import sqlite3
from contextlib import closing

import httpx
import numpy as np
from telegram import (
    Update, constants,
    BotCommand, BotCommandScopeDefault, BotCommandScopeChat,
)
from telegram.ext import (
    Application,
    CommandHandler,
    MessageHandler,
    TypeHandler,
    ContextTypes,
    filters,
)

import doc_extract  # local module: file text extractors + chunker

# ----------------------------- Config -----------------------------
TELEGRAM_TOKEN = os.environ["TELEGRAM_TOKEN"]
# Any OpenAI-compatible endpoint: Ollama's /v1, OpenAI, Groq, OpenRouter, etc.
LLM_BASE_URL = os.environ.get("LLM_BASE_URL", "http://127.0.0.1:11434/v1").rstrip("/")
LLM_API_KEY = os.environ.get("LLM_API_KEY", "ollama")
LLM_MODEL = os.environ.get("LLM_MODEL", "llama3.1")
SEARXNG_URL = os.environ.get("SEARXNG_URL", "http://127.0.0.1:8888").rstrip("/")
DB_PATH = os.environ.get(
    "DB_PATH", os.path.join(os.path.dirname(os.path.abspath(__file__)), "memory.db")
)
# Edit this file directly to change the bot's persona/instructions — re-read fresh
# on every message (see get_system_prompt()), so no restart is needed.
SYSTEM_PROMPT_PATH = os.environ.get(
    "SYSTEM_PROMPT_PATH",
    os.path.join(os.path.dirname(os.path.abspath(__file__)), "system_prompt.txt"),
)

RECENT_MESSAGES = int(os.environ.get("RECENT_MESSAGES", "12"))   # how many turns kept verbatim in context
SUMMARIZE_AFTER = int(os.environ.get("SUMMARIZE_AFTER", "24"))   # summarize once history grows past this
MAX_TOOL_ITERS = int(os.environ.get("MAX_TOOL_ITERS", "6"))      # safety cap on tool-call loops
SEARCH_RESULTS = int(os.environ.get("SEARCH_RESULTS", "5"))      # results pulled from SearXNG
HTTP_TIMEOUT = float(os.environ.get("HTTP_TIMEOUT", "120"))      # seconds per Ollama call

# Vector (long-term) memory — independently configurable since not every OpenAI-compatible
# provider offers embeddings, and embedding dimensions vary by model.
EMBED_BASE_URL = os.environ.get("EMBED_BASE_URL", LLM_BASE_URL).rstrip("/")
EMBED_API_KEY = os.environ.get("EMBED_API_KEY", LLM_API_KEY)
EMBED_MODEL = os.environ.get("EMBED_MODEL", "nomic-embed-text")
MEMORY_TOPK = int(os.environ.get("MEMORY_TOPK", "3"))           # how many past memories to recall
MEMORY_MIN_SIM = float(os.environ.get("MEMORY_MIN_SIM", "0.60"))  # cosine similarity threshold (0-1)
MEMORY_MAX_PER_CHAT = int(os.environ.get("MEMORY_MAX_PER_CHAT", "2000"))     # keep at most N newest memories per chat (0=unlimited)
# nomic-embed-text is trained with task prefixes; using them sharpens retrieval quality.
EMBED_QUERY_PREFIX = os.environ.get("EMBED_QUERY_PREFIX", "search_query: ")
EMBED_DOC_PREFIX = os.environ.get("EMBED_DOC_PREFIX", "search_document: ")

STAGING_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), ".staging")
SUPPORTED_UPLOAD = (".pdf", ".docx", ".xlsx", ".pptx", ".txt", ".md", ".png", ".jpg", ".jpeg")
MAX_UPLOAD_MB = int(os.environ.get("MAX_UPLOAD_MB", "20"))  # Telegram bot download limit ~20MB

# "Chat with file": any user can upload a doc; the bot reads it (extract+embed) and answers
# from it, scoped to that chat.
CHATDOC_MAX_INGEST = int(os.environ.get("CHATDOC_MAX_INGEST", "150"))    # max chunks embedded per upload
CHATDOC_MAX_PER_CHAT = int(os.environ.get("CHATDOC_MAX_PER_CHAT", "600"))  # cap stored chunks per chat
CHATDOC_TOPK = int(os.environ.get("CHATDOC_TOPK", "5"))
CHATDOC_MIN_SIM = float(os.environ.get("CHATDOC_MIN_SIM", "0.40"))
# Admins: higher tier for sensitive commands (e.g. /restart_bot, /whois).
ADMIN_IDS = {
    int(x) for x in os.environ.get("ADMIN_IDS", "").replace(";", ",").split(",")
    if x.strip().isdigit()
}


def is_admin(uid):
    return uid in ADMIN_IDS

# ----------------------------- Bot info (static command content) -----------------------------
# Edit these strings (or set your own persona in system_prompt.txt) to brand the bot.
START_TIME = time.time()
_DEFAULT_SKILL_TEXT = (
    "Here's what I can do 👇\n\n"
    "• 🔍 Real-time web search — current news, prices, schedules, and other up-to-date topics.\n\n"
    "• 💬 Conversation memory — I remember context from this session so you don't have to repeat yourself.\n\n"
    "• 📝 Text tasks — writing, summarizing, translation, brainstorming, editing, drafting messages.\n\n"
    "• 📄 Document Q&A — upload a file (PDF, Word, Excel, PowerPoint, txt, image) and ask about it.\n\n"
    "• 🖼️ Vision & image tools — describe photos, search images, edit images, build diagrams.\n\n"
    "• 🛠️ Network diagnostics — ping, DNS, whois, subnet calculations.\n\n"
    "• ⚡ Available 24/7 on Telegram, fast responses."
)
# /name, /owner, /skill are config-driven (set via .env) — edit .env, not this file,
# to rebrand the bot. BOT_SKILL supports \n for line breaks (e.g. "Line one\nLine two").
INFO = {
    "name": os.environ.get("BOT_NAME", "🤖 *Telegram Bot Kit* — a generic AI assistant bot."),
    "owner": os.environ.get("BOT_OWNER", "👤 Owner: (set BOT_OWNER in .env)"),
    "skill": os.environ.get("BOT_SKILL", _DEFAULT_SKILL_TEXT).replace("\\n", "\n"),
}

# Fallback only — the live prompt is loaded from SYSTEM_PROMPT_PATH by get_system_prompt().
# This keeps the bot working even if that file is ever missing or unreadable.
_DEFAULT_SYSTEM_PROMPT = (
    "You are a helpful AI assistant on Telegram. Answer clearly and honestly, "
    "in the same language the user writes in."
)


def get_system_prompt():
    """Read the external prompt file fresh each call (hot-reload, no restart
    needed to pick up edits). Falls back to the built-in default if missing."""
    try:
        with open(SYSTEM_PROMPT_PATH, "r", encoding="utf-8") as f:
            text = f.read().strip()
        if text:
            return text
    except OSError as e:
        log.warning("could not read %s, using default system prompt: %s", SYSTEM_PROMPT_PATH, e)
    return _DEFAULT_SYSTEM_PROMPT


TOOLS = [
    {
        "type": "function",
        "function": {
            "name": "web_search",
            "description": (
                "Cari informasi terkini di web. Gunakan untuk berita, harga, cuaca, "
                "kejadian terbaru, atau fakta yang bisa berubah seiring waktu."
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "query": {"type": "string", "description": "kata kunci pencarian"}
                },
                "required": ["query"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "web_image_search",
            "description": (
                "Cari GAMBAR di web (diagram, foto, ilustrasi, arsitektur jaringan, screenshot, dll) "
                "lalu tampilkan ke user. Gunakan saat user meminta gambar/diagram/ilustrasi/visual."
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "query": {"type": "string", "description": "kata kunci gambar"}
                },
                "required": ["query"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "buat_diagram",
            "description": (
                "Buat diagram topologi jaringan/arsitektur (Graphviz DOT) atau flowchart/alur proses "
                "(Mermaid) dari deskripsi user, lalu kirim sebagai gambar. Gunakan ini (bukan "
                "web_image_search) saat user minta DIBUATKAN diagram spesifik sesuai kebutuhan mereka "
                "(topologi jaringan mereka, alur proses tertentu) — bukan mencari gambar contoh yang "
                "sudah ada di internet."
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "tipe": {
                        "type": "string",
                        "enum": ["graphviz", "mermaid"],
                        "description": (
                            "graphviz untuk topologi/diagram jaringan/arsitektur (node & koneksi), "
                            "mermaid untuk flowchart/alur proses/urutan langkah."
                        ),
                    },
                    "kode": {
                        "type": "string",
                        "description": (
                            "Kode diagram valid. Untuk graphviz: 'digraph G { A -> B }'. "
                            "Untuk mermaid: 'graph TD; A-->B;'"
                        ),
                    },
                },
                "required": ["tipe", "kode"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "buat_dokumen",
            "description": (
                "Buat file dokumen (txt, docx, pptx, pdf, md, csv, json, atau html) dari konten yang kamu "
                "susun sendiri, lalu kirim ke user sebagai file. Gunakan saat user minta dibuatkan laporan/"
                "ringkasan/proposal/dokumentasi/data dalam bentuk FILE yang bisa diunduh/diedit — bukan "
                "cuma jawaban di chat biasa."
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "tipe": {
                        "type": "string",
                        "enum": ["txt", "docx", "pptx", "pdf", "md", "csv", "json", "html"],
                    },
                    "judul": {"type": "string", "description": "judul dokumen"},
                    "konten": {
                        "type": "string",
                        "description": (
                            "Isi dokumen. Formatnya beda tergantung tipe:\n"
                            "- txt/docx/pptx/pdf/md: format sederhana — baris diawali '# ' = heading utama, "
                            "'## ' = sub-heading, '- ' = bullet point, baris kosong = pisah paragraf. "
                            "Untuk pptx, tiap heading ('# '/'## ') memulai slide baru.\n"
                            "- csv: tulis langsung sebagai CSV valid (baris pertama = header, kolom "
                            "dipisah koma).\n"
                            "- json: tulis langsung sebagai JSON valid.\n"
                            "- html: tulis langsung sebagai HTML (boleh cuma fragment body, tidak perlu "
                            "<html>/<head> lengkap)."
                        ),
                    },
                },
                "required": ["tipe", "judul", "konten"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "edit_gambar",
            "description": (
                "Edit gambar TERAKHIR di chat ini (foto yang baru diupload, diagram yang baru dibuat, "
                "atau hasil edit sebelumnya) — crop, resize, rotate, atau tempel teks/watermark. Ini "
                "edit TEKNIS/MEKANIS saja (bukan mengubah isi/konten gambar secara AI)."
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "aksi": {"type": "string", "enum": ["crop", "resize", "rotate", "teks"]},
                    "resize_lebar": {"type": "integer", "description": "lebar baru dalam pixel (aksi resize)"},
                    "resize_tinggi": {
                        "type": "integer",
                        "description": "tinggi baru dalam pixel (aksi resize, opsional — proporsional kalau kosong)",
                    },
                    "rotate_derajat": {
                        "type": "number",
                        "description": "derajat rotasi searah jarum jam (aksi rotate)",
                    },
                    "crop_kiri": {"type": "number", "description": "persen dipotong dari kiri, 0-100 (aksi crop)"},
                    "crop_atas": {"type": "number", "description": "persen dipotong dari atas, 0-100 (aksi crop)"},
                    "crop_kanan": {"type": "number", "description": "persen dipotong dari kanan, 0-100 (aksi crop)"},
                    "crop_bawah": {"type": "number", "description": "persen dipotong dari bawah, 0-100 (aksi crop)"},
                    "teks": {"type": "string", "description": "teks watermark/anotasi (aksi teks)"},
                    "posisi": {
                        "type": "string",
                        "enum": ["atas-kiri", "atas-kanan", "bawah-kiri", "bawah-kanan", "tengah"],
                        "description": "posisi teks (aksi teks), default bawah-kanan",
                    },
                },
                "required": ["aksi"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "network_tool",
            "description": (
                "Alat diagnostik jaringan: ping, dns (resolve hostname), whois, atau subnet "
                "(hitung network/broadcast/range dari CIDR). Read-only, hasil langsung ditampilkan."
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "aksi": {"type": "string", "enum": ["ping", "dns", "whois", "subnet"]},
                    "target": {
                        "type": "string",
                        "description": "hostname/IP (ping/dns/whois) atau CIDR mis. '10.0.0.0/24' (subnet)",
                    },
                },
                "required": ["aksi", "target"],
            },
        },
    },
]

KROKI_URL = os.environ.get("KROKI_URL", "http://127.0.0.1:8001").rstrip("/")
DOC_TYPES = ("txt", "docx", "pptx", "pdf", "md", "csv", "json", "html")
DOC_EXTS = tuple(f".{t}" for t in DOC_TYPES)


def _parse_doc_blocks(konten):
    """Split a lightweight-markdown content string into (kind, text) blocks."""
    blocks = []
    for line in (konten or "").splitlines():
        s = line.strip()
        if not s:
            continue
        if s.startswith("## "):
            blocks.append(("h2", s[3:].strip()))
        elif s.startswith("# "):
            blocks.append(("h1", s[2:].strip()))
        elif s.startswith(("- ", "* ")):
            blocks.append(("bullet", s[2:].strip()))
        else:
            blocks.append(("p", s))
    return blocks


def _build_txt(judul, konten, path):
    with open(path, "w", encoding="utf-8") as f:
        f.write(f"{judul}\n{'=' * len(judul)}\n\n{konten.strip()}\n")


def _build_docx(judul, konten, path):
    from docx import Document

    doc = Document()
    doc.add_heading(judul, level=0)
    for kind, text in _parse_doc_blocks(konten):
        if kind == "h1":
            doc.add_heading(text, level=1)
        elif kind == "h2":
            doc.add_heading(text, level=2)
        elif kind == "bullet":
            doc.add_paragraph(text, style="List Bullet")
        else:
            doc.add_paragraph(text)
    doc.save(path)


def _group_slides(judul, konten):
    slides = []
    cur = None
    for kind, text in _parse_doc_blocks(konten):
        if kind in ("h1", "h2"):
            cur = {"title": text, "lines": []}
            slides.append(cur)
        else:
            if cur is None:
                cur = {"title": judul, "lines": []}
                slides.append(cur)
            cur["lines"].append(text)
    return slides or [{"title": judul, "lines": []}]


def _build_pptx(judul, konten, path):
    from pptx import Presentation

    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[0]).shapes.title.text = judul
    for sl in _group_slides(judul, konten):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = sl["title"]
        tf = slide.placeholders[1].text_frame
        tf.clear()
        if sl["lines"]:
            tf.text = sl["lines"][0]
            for line in sl["lines"][1:]:
                tf.add_paragraph().text = line
    prs.save(path)


def _build_pdf(judul, konten, path):
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer

    styles = getSampleStyleSheet()
    story = [Paragraph(judul, styles["Title"]), Spacer(1, 12)]
    for kind, text in _parse_doc_blocks(konten):
        style = {"h1": "Heading1", "h2": "Heading2"}.get(kind, "Normal")
        text_out = f"&bull; {text}" if kind == "bullet" else text
        story.append(Paragraph(text_out, styles[style]))
        story.append(Spacer(1, 6))
    SimpleDocTemplate(path, pagesize=A4).build(story)


def _build_md(judul, konten, path):
    with open(path, "w", encoding="utf-8") as f:
        f.write(f"# {judul}\n\n{konten.strip()}\n")


def _build_csv(judul, konten, path):
    import csv
    import io

    rows = list(csv.reader(io.StringIO(konten.strip())))
    if not rows:
        raise ValueError("konten CSV kosong")
    with open(path, "w", encoding="utf-8", newline="") as f:
        csv.writer(f).writerows(rows)


def _build_json(judul, konten, path):
    parsed = json.loads(konten)  # raises ValueError if invalid -> caller reports back to the model
    with open(path, "w", encoding="utf-8") as f:
        json.dump(parsed, f, indent=2, ensure_ascii=False)


def _build_html(judul, konten, path):
    body = konten.strip()
    if "<html" not in body.lower():
        body = (
            f"<!DOCTYPE html><html><head><meta charset=\"utf-8\">"
            f"<title>{judul}</title></head><body><h1>{judul}</h1>\n{body}\n</body></html>"
        )
    with open(path, "w", encoding="utf-8") as f:
        f.write(body)


def build_document(tipe, judul, konten, out_dir):
    """Build a document file of the given type. Returns the local file path."""
    os.makedirs(out_dir, exist_ok=True)
    safe_name = re.sub(r"[^\w\-]+", "_", judul).strip("_")[:50] or "dokumen"
    path = os.path.join(out_dir, f"{safe_name}_{int(time.time())}_{uuid.uuid4().hex[:6]}.{tipe}")
    builders = {
        "txt": _build_txt, "docx": _build_docx, "pptx": _build_pptx, "pdf": _build_pdf,
        "md": _build_md, "csv": _build_csv, "json": _build_json, "html": _build_html,
    }
    builders[tipe](judul, konten, path)
    return path


# ----------------------------- Image editing (mechanical, not AI) -----------------------------
LAST_IMAGE = {}  # chat_id -> local path of the most recent bot-owned image (upload/diagram/edit)


def register_last_image(chat_id, path):
    LAST_IMAGE[chat_id] = path


def edit_image(chat_id, aksi, resize_lebar=None, resize_tinggi=None, rotate_derajat=None,
                crop_kiri=None, crop_atas=None, crop_kanan=None, crop_bawah=None,
                teks=None, posisi="bawah-kanan"):
    """Apply a mechanical edit (crop/resize/rotate/text-overlay) to the chat's last
    bot-owned image. Returns the new local PNG path."""
    from PIL import Image, ImageDraw, ImageFont

    src = LAST_IMAGE.get(chat_id)
    if not src or not os.path.isfile(src):
        raise ValueError("Belum ada gambar di chat ini yang bisa diedit (upload foto atau buat diagram dulu).")

    img = Image.open(src).convert("RGBA")
    if aksi == "resize":
        if not resize_lebar:
            raise ValueError("resize_lebar wajib diisi.")
        w = int(resize_lebar)
        h = int(resize_tinggi) if resize_tinggi else int(img.height * (w / img.width))
        img = img.resize((w, h))
    elif aksi == "rotate":
        img = img.rotate(-(rotate_derajat or 0), expand=True)
    elif aksi == "crop":
        w, h = img.size
        left = int(w * (crop_kiri or 0) / 100)
        top = int(h * (crop_atas or 0) / 100)
        right = int(w * (1 - (crop_kanan or 0) / 100))
        bottom = int(h * (1 - (crop_bawah or 0) / 100))
        if right <= left or bottom <= top:
            raise ValueError("Persen crop terlalu besar, hasil crop-nya jadi kosong.")
        img = img.crop((left, top, right, bottom))
    elif aksi == "teks":
        if not teks:
            raise ValueError("teks wajib diisi.")
        draw = ImageDraw.Draw(img)
        try:
            font = ImageFont.truetype("DejaVuSans-Bold.ttf", max(16, img.width // 30))
        except Exception:  # noqa: BLE001
            font = ImageFont.load_default()
        bbox = draw.textbbox((0, 0), teks, font=font)
        tw, th = bbox[2] - bbox[0], bbox[3] - bbox[1]
        pad = 8
        pos_map = {
            "atas-kiri": (pad, pad),
            "atas-kanan": (img.width - tw - pad * 3, pad),
            "bawah-kiri": (pad, img.height - th - pad * 3),
            "bawah-kanan": (img.width - tw - pad * 3, img.height - th - pad * 3),
            "tengah": ((img.width - tw) // 2, (img.height - th) // 2),
        }
        x, y = pos_map.get(posisi, pos_map["bawah-kanan"])
        draw.rectangle([x - pad, y - pad, x + tw + pad, y + th + pad], fill=(0, 0, 0, 160))
        draw.text((x, y), teks, fill=(255, 255, 255, 255), font=font)
    else:
        raise ValueError(f"aksi tidak dikenal: {aksi}")

    os.makedirs(STAGING_DIR, exist_ok=True)
    out_path = os.path.join(STAGING_DIR, f"edit_{chat_id}_{int(time.time())}_{uuid.uuid4().hex[:6]}.png")
    img.convert("RGB").save(out_path, "PNG")
    LAST_IMAGE[chat_id] = out_path
    return out_path


# ----------------------------- network utility tools -----------------------------
NET_TARGET_RE = re.compile(r"^[a-zA-Z0-9.\-:_/]+$")


def network_tool(aksi, target):
    """ping / dns / whois / subnet — read-only diagnostics, no external paid API.
    Target is validated against a strict charset before ever reaching subprocess."""
    target = (target or "").strip()
    if not target:
        raise ValueError("target kosong.")

    if aksi == "subnet":
        try:
            net = ipaddress.ip_network(target, strict=False)
        except ValueError as e:
            raise ValueError(f"CIDR/IP tidak valid: {e}") from e
        broadcast = str(net.broadcast_address) if net.version == 4 else "-"
        return (
            f"Network   : {net.network_address}\n"
            f"Broadcast : {broadcast}\n"
            f"Netmask   : {net.netmask}\n"
            f"Prefix    : /{net.prefixlen}\n"
            f"Jumlah IP : {net.num_addresses}\n"
            f"Range     : {net[0]} - {net[-1]}"
        )

    if not NET_TARGET_RE.match(target):
        raise ValueError("target mengandung karakter yang tidak diizinkan.")

    if aksi == "ping":
        r = subprocess.run(
            ["ping", "-c", "4", "-W", "2", target], capture_output=True, text=True, timeout=15
        )
        return (r.stdout or r.stderr or "Tidak ada output.").strip()[:2500]

    if aksi == "dns":
        r = subprocess.run(["dig", "+short", target], capture_output=True, text=True, timeout=10)
        out = r.stdout.strip()
        return out or "Tidak ada hasil DNS (kemungkinan domain tidak resolve)."

    if aksi == "whois":
        r = subprocess.run(["whois", target], capture_output=True, text=True, timeout=15)
        return (r.stdout or r.stderr or "Tidak ada hasil whois.").strip()[:3000]

    raise ValueError(f"aksi tidak dikenal: {aksi}")


async def kroki_render(tipe, kode):
    """Render diagram source to PNG bytes via a self-hosted Kroki instance.
    Returns (png_bytes, None) on success or (None, error_message) on failure."""
    url = f"{KROKI_URL}/{tipe}/png"
    try:
        async with httpx.AsyncClient(timeout=30) as client:
            resp = await client.post(url, content=kode.encode("utf-8"),
                                      headers={"Content-Type": "text/plain"})
        if resp.status_code != 200:
            return None, f"Kroki error {resp.status_code}: {resp.text[:200]}"
        return resp.content, None
    except Exception as e:  # noqa: BLE001
        return None, str(e)


logging.basicConfig(
    format="%(asctime)s %(levelname)s %(name)s: %(message)s", level=logging.INFO
)
log = logging.getLogger("telegram-bot-kit")


# ----------------------------- Storage (SQLite) -----------------------------
def db_init():
    with closing(sqlite3.connect(DB_PATH, timeout=30)) as con:
        con.execute("PRAGMA journal_mode=WAL")
        con.execute(
            """CREATE TABLE IF NOT EXISTS messages(
                   id INTEGER PRIMARY KEY AUTOINCREMENT,
                   chat_id INTEGER NOT NULL,
                   role    TEXT    NOT NULL,
                   content TEXT    NOT NULL,
                   ts      REAL    NOT NULL)"""
        )
        con.execute(
            """CREATE TABLE IF NOT EXISTS summaries(
                   chat_id INTEGER PRIMARY KEY,
                   summary TEXT NOT NULL,
                   updated REAL NOT NULL)"""
        )
        con.execute(
            """CREATE TABLE IF NOT EXISTS memories(
                   id        INTEGER PRIMARY KEY AUTOINCREMENT,
                   chat_id   INTEGER NOT NULL,
                   text      TEXT    NOT NULL,
                   embedding BLOB    NOT NULL,
                   ts        REAL    NOT NULL)"""
        )
        con.execute(
            """CREATE TABLE IF NOT EXISTS chat_docs(
                   id        INTEGER PRIMARY KEY AUTOINCREMENT,
                   chat_id   INTEGER NOT NULL,
                   name      TEXT    NOT NULL,
                   chunk_idx INTEGER,
                   text      TEXT    NOT NULL,
                   embedding BLOB    NOT NULL,
                   ts        REAL    NOT NULL)"""
        )
        con.execute("CREATE INDEX IF NOT EXISTS idx_msg_chat ON messages(chat_id, id)")
        con.execute(
            """CREATE TABLE IF NOT EXISTS seen_users(
                   telegram_id INTEGER PRIMARY KEY,
                   name        TEXT,
                   username    TEXT,
                   first_seen  REAL,
                   last_seen   REAL,
                   msg_count   INTEGER DEFAULT 0)"""
        )
        con.execute("CREATE INDEX IF NOT EXISTS idx_chatdoc ON chat_docs(chat_id, id)")
        con.execute("CREATE INDEX IF NOT EXISTS idx_mem_chat ON memories(chat_id)")
        con.commit()


def db_add_message(chat_id, role, content):
    with closing(sqlite3.connect(DB_PATH, timeout=30)) as con:
        con.execute(
            "INSERT INTO messages(chat_id, role, content, ts) VALUES (?,?,?,?)",
            (chat_id, role, content, time.time()),
        )
        con.commit()


def db_recent(chat_id, limit):
    with closing(sqlite3.connect(DB_PATH, timeout=30)) as con:
        rows = con.execute(
            "SELECT role, content FROM messages WHERE chat_id=? ORDER BY id DESC LIMIT ?",
            (chat_id, limit),
        ).fetchall()
    return [{"role": r, "content": c} for r, c in reversed(rows)]


def db_count(chat_id):
    with closing(sqlite3.connect(DB_PATH, timeout=30)) as con:
        return con.execute(
            "SELECT COUNT(*) FROM messages WHERE chat_id=?", (chat_id,)
        ).fetchone()[0]


def db_get_summary(chat_id):
    with closing(sqlite3.connect(DB_PATH, timeout=30)) as con:
        row = con.execute(
            "SELECT summary FROM summaries WHERE chat_id=?", (chat_id,)
        ).fetchone()
    return row[0] if row else ""


def db_set_summary(chat_id, summary):
    with closing(sqlite3.connect(DB_PATH, timeout=30)) as con:
        con.execute(
            """INSERT INTO summaries(chat_id, summary, updated) VALUES(?,?,?)
               ON CONFLICT(chat_id) DO UPDATE SET
                   summary=excluded.summary, updated=excluded.updated""",
            (chat_id, summary, time.time()),
        )
        con.commit()


def db_prune_older(chat_id, keep):
    """Delete all but the most recent `keep` messages for a chat."""
    with closing(sqlite3.connect(DB_PATH, timeout=30)) as con:
        con.execute(
            """DELETE FROM messages
               WHERE chat_id=? AND id NOT IN (
                   SELECT id FROM messages WHERE chat_id=? ORDER BY id DESC LIMIT ?)""",
            (chat_id, chat_id, keep),
        )
        con.commit()


def db_reset(chat_id):
    with closing(sqlite3.connect(DB_PATH, timeout=30)) as con:
        con.execute("DELETE FROM messages WHERE chat_id=?", (chat_id,))
        con.execute("DELETE FROM summaries WHERE chat_id=?", (chat_id,))
        con.execute("DELETE FROM memories WHERE chat_id=?", (chat_id,))
        con.execute("DELETE FROM chat_docs WHERE chat_id=?", (chat_id,))
        con.commit()


def db_touch_user(uid, name, username):
    """Record that a user interacted with the bot (for the /users seen-list)."""
    now = time.time()
    with closing(sqlite3.connect(DB_PATH, timeout=30)) as con:
        con.execute(
            """INSERT INTO seen_users(telegram_id, name, username, first_seen, last_seen, msg_count)
               VALUES(?,?,?,?,?,1)
               ON CONFLICT(telegram_id) DO UPDATE SET
                   name=COALESCE(excluded.name, seen_users.name),
                   username=COALESCE(excluded.username, seen_users.username),
                   last_seen=excluded.last_seen,
                   msg_count=seen_users.msg_count + 1""",
            (uid, name, username, now, now),
        )
        con.commit()


def db_list_users(limit=100):
    with closing(sqlite3.connect(DB_PATH, timeout=30)) as con:
        return con.execute(
            "SELECT telegram_id, name, username, last_seen, msg_count FROM seen_users "
            "ORDER BY last_seen DESC LIMIT ?",
            (limit,),
        ).fetchall()


def db_add_memory(chat_id, text, emb):
    with closing(sqlite3.connect(DB_PATH, timeout=30)) as con:
        con.execute(
            "INSERT INTO memories(chat_id, text, embedding, ts) VALUES (?,?,?,?)",
            (chat_id, text, np.asarray(emb, dtype=np.float32).tobytes(), time.time()),
        )
        # Size cap: keep only the newest MEMORY_MAX_PER_CHAT memories per chat.
        if MEMORY_MAX_PER_CHAT > 0:
            con.execute(
                """DELETE FROM memories WHERE chat_id=? AND id NOT IN (
                       SELECT id FROM memories WHERE chat_id=? ORDER BY id DESC LIMIT ?)""",
                (chat_id, chat_id, MEMORY_MAX_PER_CHAT),
            )
        con.commit()


def db_search_memory(chat_id, query_emb, k=MEMORY_TOPK, min_sim=MEMORY_MIN_SIM, exclude_texts=()):
    """Return up to k stored memory texts most similar (cosine) to query_emb."""
    with closing(sqlite3.connect(DB_PATH, timeout=30)) as con:
        rows = con.execute(
            "SELECT text, embedding FROM memories WHERE chat_id=?", (chat_id,)
        ).fetchall()
    if not rows:
        return []
    q = np.asarray(query_emb, dtype=np.float32)
    qn = q / (np.linalg.norm(q) + 1e-8)
    exclude = set(exclude_texts)
    scored = []
    for text, blob in rows:
        if text in exclude:
            continue
        v = np.frombuffer(blob, dtype=np.float32)
        if v.shape != q.shape:
            continue
        sim = float(np.dot(qn, v / (np.linalg.norm(v) + 1e-8)))
        if sim >= min_sim:
            scored.append((sim, text))
    scored.sort(reverse=True)
    return [t for _, t in scored[:k]]


def db_add_chat_doc(chat_id, name, rows):
    """Store an uploaded doc's chunks for a chat. rows = [(chunk_idx, text, emb_bytes), ...].
    Caps the per-chat chunk count so storage stays light."""
    if not rows:
        return
    now = time.time()
    with closing(sqlite3.connect(DB_PATH, timeout=30)) as con:
        con.executemany(
            "INSERT INTO chat_docs(chat_id, name, chunk_idx, text, embedding, ts) VALUES (?,?,?,?,?,?)",
            [(chat_id, name, idx, txt, blob, now) for idx, txt, blob in rows],
        )
        if CHATDOC_MAX_PER_CHAT > 0:
            con.execute(
                """DELETE FROM chat_docs WHERE chat_id=? AND id NOT IN (
                       SELECT id FROM chat_docs WHERE chat_id=? ORDER BY id DESC LIMIT ?)""",
                (chat_id, chat_id, CHATDOC_MAX_PER_CHAT),
            )
        con.commit()


def db_latest_chat_doc_name(chat_id):
    with closing(sqlite3.connect(DB_PATH, timeout=30)) as con:
        row = con.execute(
            "SELECT name FROM chat_docs WHERE chat_id=? ORDER BY id DESC LIMIT 1", (chat_id,)
        ).fetchone()
    return row[0] if row else None


def db_chat_doc_head(chat_id, name=None, limit=15):
    """First `limit` chunks of an uploaded doc (latest one if name omitted) — good for summaries."""
    with closing(sqlite3.connect(DB_PATH, timeout=30)) as con:
        if name is None:
            row = con.execute(
                "SELECT name FROM chat_docs WHERE chat_id=? ORDER BY id DESC LIMIT 1", (chat_id,)
            ).fetchone()
            if not row:
                return []
            name = row[0]
        rows = con.execute(
            "SELECT text, name FROM chat_docs WHERE chat_id=? AND name=? ORDER BY chunk_idx LIMIT ?",
            (chat_id, name, limit),
        ).fetchall()
    return [(t, n) for t, n in rows]


def db_search_chat_docs(chat_id, query_emb, k=CHATDOC_TOPK, min_sim=CHATDOC_MIN_SIM):
    """Cosine search over this chat's uploaded docs; return up to k (text, name)."""
    with closing(sqlite3.connect(DB_PATH, timeout=30)) as con:
        rows = con.execute(
            "SELECT text, name, embedding FROM chat_docs WHERE chat_id=?", (chat_id,)
        ).fetchall()
    if not rows:
        return []
    q = np.asarray(query_emb, dtype=np.float32)
    qn = q / (np.linalg.norm(q) + 1e-8)
    scored = []
    for text, name, blob in rows:
        v = np.frombuffer(blob, dtype=np.float32)
        if v.shape != q.shape:
            continue
        sim = float(np.dot(qn, v / (np.linalg.norm(v) + 1e-8)))
        if sim >= min_sim:
            scored.append((sim, text, name))
    scored.sort(reverse=True)
    return [(t, n) for _, t, n in scored[:k]]




# ----------------------------- External calls -----------------------------
async def searxng_search(query, n=SEARCH_RESULTS):
    """Query SearXNG JSON API and return a compact text block of results."""
    try:
        async with httpx.AsyncClient(timeout=30) as client:
            r = await client.get(
                f"{SEARXNG_URL}/search", params={"q": query, "format": "json"}
            )
            r.raise_for_status()
            data = r.json()
    except Exception as e:  # noqa: BLE001
        log.warning("searxng error: %s", e)
        return f"(pencarian gagal: {e})"

    results = data.get("results", [])[:n]
    if not results:
        return "Tidak ada hasil yang ditemukan."
    lines = []
    for i, res in enumerate(results, 1):
        title = res.get("title", "").strip()
        url = res.get("url", "").strip()
        content = (res.get("content") or "").strip()
        lines.append(f"{i}. {title}\n   {url}\n   {content}")
    return "\n".join(lines)


async def searxng_image_search(query, n=5):
    """Query SearXNG image search; return up to n (img_url, title) tuples."""
    try:
        async with httpx.AsyncClient(timeout=30) as client:
            r = await client.get(
                f"{SEARXNG_URL}/search",
                params={"q": query, "format": "json", "categories": "images"},
            )
            r.raise_for_status()
            data = r.json()
    except Exception as e:  # noqa: BLE001
        log.warning("image search error: %s", e)
        return []
    bad = ("jsdelivr.net", "lucide", "devicon", "/icons/", "icon-")
    out = []
    for res in data.get("results", []):
        src = res.get("img_src") or res.get("thumbnail_src") or ""
        if not src.startswith(("http://", "https://")):
            continue
        low = src.lower()
        if low.endswith(".svg") or any(b in low for b in bad):
            continue  # skip UI icons / vector junk
        out.append((src, (res.get("title") or "").strip()[:80]))
        if len(out) >= n:
            break
    return out


async def llm_chat(messages, tools=None):
    """Call any OpenAI-compatible /chat/completions endpoint. Returns the raw JSON
    response; use choices[0]["message"] to get the assistant turn."""
    payload = {"model": LLM_MODEL, "messages": messages, "stream": False}
    if tools:
        payload["tools"] = tools
    headers = {"Authorization": f"Bearer {LLM_API_KEY}"}
    async with httpx.AsyncClient(timeout=HTTP_TIMEOUT) as client:
        r = await client.post(f"{LLM_BASE_URL}/chat/completions", json=payload, headers=headers)
        r.raise_for_status()
        return r.json()


async def embed(text, task=None):
    """Return a float32 embedding for text, or None on failure (graceful fallback).

    task="query"  -> prefix for a search query
    task="document" -> prefix for stored content
    """
    if task == "query":
        text = EMBED_QUERY_PREFIX + text
    elif task == "document":
        text = EMBED_DOC_PREFIX + text
    try:
        headers = {"Authorization": f"Bearer {EMBED_API_KEY}"}
        async with httpx.AsyncClient(timeout=30) as client:
            r = await client.post(
                f"{EMBED_BASE_URL}/embeddings",
                json={"model": EMBED_MODEL, "input": text},
                headers=headers,
            )
            r.raise_for_status()
            data = r.json().get("data") or []
            vec = data[0].get("embedding") if data else None
        if not vec:
            return None
        return np.asarray(vec, dtype=np.float32)
    except Exception as e:  # noqa: BLE001
        log.warning("embed failed: %s", e)
        return None


# ----------------------------- Core logic -----------------------------
async def generate_reply(chat_id, user_text, images=None):
    summary = db_get_summary(chat_id)
    history = db_recent(chat_id, RECENT_MESSAGES)

    # Vector (long-term) memory + per-chat uploaded documents: recall by semantic similarity.
    query_emb = await embed(user_text, task="query")
    recalled = []
    chatdoc_hits = []
    if query_emb is not None:
        recent_texts = {m["content"] for m in history}
        recalled = db_search_memory(chat_id, query_emb, exclude_texts=recent_texts)

        # Uploaded-file RAG ("chat with file") — available to EVERYONE, scoped to this chat.
        if db_latest_chat_doc_name(chat_id):
            ql2 = user_text.lower()
            if re.search(
                r"(rangkum|ringkas|ringkasan|summary|\bisi\b|isinya|jelas(kan|in)|poin penting|apa isi)",
                ql2,
            ):
                chatdoc_hits = db_chat_doc_head(chat_id, limit=15)
            else:
                chatdoc_hits = db_search_chat_docs(chat_id, query_emb)

    # Media to attach to the reply: web image search results are added later by the tool handler.
    media = []

    messages = [{"role": "system", "content": get_system_prompt()}]
    if summary:
        messages.append(
            {"role": "system", "content": f"Ringkasan percakapan sebelumnya:\n{summary}"}
        )
    if recalled:
        block = "\n".join(f"- {t}" for t in recalled)
        messages.append(
            {
                "role": "system",
                "content": (
                    "Ingatan relevan dari percakapan lampau (pakai kalau membantu, "
                    "abaikan kalau tidak relevan):\n" + block
                ),
            }
        )
    messages.extend(history)
    # Put the uploaded-file content LAST (highest salience) so it overrides any earlier
    # turn where the bot wrongly claimed it couldn't read files.
    if chatdoc_hits:
        doc_block = "\n\n".join(f"[Dokumen: {name}]\n{txt}" for txt, name in chatdoc_hits)
        messages.append(
            {
                "role": "system",
                "content": (
                    "PENTING: user mengunggah file di chat ini dan ISINYA TERSEDIA di bawah — kamu "
                    "BISA membacanya. Abaikan pernyataanmu sebelumnya (jika ada) yang bilang tidak "
                    "bisa membaca/akses file; itu keliru. Jawab/rangkum berdasarkan isi ini dan sebut "
                    "nama filenya:\n\n" + doc_block
                ),
            }
        )
    user_msg = {"role": "user", "content": user_text}
    if images:
        user_msg["images"] = images
    messages.append(user_msg)

    db_add_message(chat_id, "user", user_text)
    doc_emb = await embed(user_text, task="document")
    if doc_emb is not None:
        db_add_memory(chat_id, user_text, doc_emb)

    for _ in range(MAX_TOOL_ITERS):
        data = await llm_chat(messages, tools=TOOLS)
        msg = (data.get("choices") or [{}])[0].get("message", {}) or {}
        tool_calls = msg.get("tool_calls")

        if tool_calls:
            # Echo back the assistant turn that requested the tools, then the results.
            messages.append(
                {
                    "role": "assistant",
                    "content": msg.get("content", "") or "",
                    "tool_calls": tool_calls,
                }
            )
            for tc in tool_calls:
                fn = tc.get("function", {}) or {}
                name = fn.get("name")
                args = fn.get("arguments", {})
                if isinstance(args, str):
                    try:
                        args = json.loads(args)
                    except Exception:  # noqa: BLE001
                        args = {}
                if name == "web_search":
                    query = args.get("query", "")
                    log.info("chat %s web_search: %s", chat_id, query)
                    result = await searxng_search(query)
                elif name == "web_image_search":
                    query = args.get("query", "")
                    log.info("chat %s web_image_search: %s", chat_id, query)
                    imgs = await searxng_image_search(query)
                    for url, title in imgs:
                        media.append((url, f"🌐 {title}" if title else "🌐 gambar"))
                    if imgs:
                        result = (
                            f"{len(imgs)} gambar untuk '{query}' SUDAH dikirim ke user. "
                            "Beri keterangan singkat saja; JANGAN menempelkan URL gambar di teks jawaban."
                        )
                    else:
                        result = "Tidak ada gambar yang ditemukan."
                elif name == "buat_diagram":
                    tipe = args.get("tipe", "")
                    kode = args.get("kode", "")
                    log.info("chat %s buat_diagram: %s", chat_id, tipe)
                    if tipe not in ("graphviz", "mermaid"):
                        result = "Tipe diagram tidak valid, harus 'graphviz' atau 'mermaid'."
                    else:
                        png, err = await kroki_render(tipe, kode)
                        if err:
                            result = f"Gagal membuat diagram: {err}. Coba perbaiki kodenya."
                        else:
                            os.makedirs(STAGING_DIR, exist_ok=True)
                            path = os.path.join(
                                STAGING_DIR, f"diagram_{chat_id}_{int(time.time())}_{uuid.uuid4().hex[:6]}.png"
                            )
                            with open(path, "wb") as f:
                                f.write(png)
                            media.append((path, "📊 diagram"))
                            register_last_image(chat_id, path)
                            result = (
                                "Diagram SUDAH dibuat dan dikirim ke user sebagai gambar. "
                                "Beri keterangan singkat saja; JANGAN menempelkan kode diagram di teks jawaban."
                            )
                elif name == "buat_dokumen":
                    tipe = args.get("tipe", "")
                    judul = args.get("judul", "Dokumen")
                    konten = args.get("konten", "")
                    log.info("chat %s buat_dokumen: %s / %s", chat_id, tipe, judul)
                    if tipe not in DOC_TYPES:
                        result = "Tipe dokumen tidak valid, harus txt/docx/pptx/pdf."
                    else:
                        try:
                            path = await asyncio.to_thread(
                                build_document, tipe, judul, konten, STAGING_DIR
                            )
                            media.append((path, f"📄 {judul}.{tipe}"))
                            result = (
                                f"Dokumen '{judul}.{tipe}' SUDAH dibuat dan dikirim ke user sebagai file. "
                                "Beri keterangan singkat saja; JANGAN menempelkan isi dokumen mentah lagi di teks jawaban."
                            )
                        except Exception as e:  # noqa: BLE001
                            log.exception("buat_dokumen failed")
                            result = f"Gagal membuat dokumen: {e}"
                elif name == "edit_gambar":
                    aksi = args.get("aksi", "")
                    log.info("chat %s edit_gambar: %s", chat_id, aksi)
                    try:
                        path = await asyncio.to_thread(
                            edit_image, chat_id, aksi,
                            resize_lebar=args.get("resize_lebar"),
                            resize_tinggi=args.get("resize_tinggi"),
                            rotate_derajat=args.get("rotate_derajat"),
                            crop_kiri=args.get("crop_kiri"), crop_atas=args.get("crop_atas"),
                            crop_kanan=args.get("crop_kanan"), crop_bawah=args.get("crop_bawah"),
                            teks=args.get("teks"), posisi=args.get("posisi", "bawah-kanan"),
                        )
                        media.append((path, "🖼️ gambar diedit"))
                        result = (
                            "Gambar SUDAH diedit dan dikirim ke user sebagai foto baru. "
                            "Beri keterangan singkat saja."
                        )
                    except Exception as e:  # noqa: BLE001
                        result = f"Gagal mengedit gambar: {e}"
                elif name == "network_tool":
                    aksi = args.get("aksi", "")
                    target = args.get("target", "")
                    log.info("chat %s network_tool: %s %s", chat_id, aksi, target)
                    try:
                        result = await asyncio.to_thread(network_tool, aksi, target)
                    except Exception as e:  # noqa: BLE001
                        result = f"Gagal menjalankan {aksi}: {e}"
                else:
                    result = f"(tool tidak dikenal: {name})"
                messages.append(
                    {
                        "role": "tool",
                        "tool_call_id": tc.get("id", ""),
                        "content": result,
                    }
                )
            continue

        # No tool call -> final answer.
        content = (msg.get("content") or "").strip()
        if content:
            db_add_message(chat_id, "assistant", content)
            await maybe_summarize(chat_id)
        return (content or "(maaf, jawaban kosong)", media[:6])

    fallback = "Maaf, terlalu banyak langkah pencarian. Coba pertanyaan yang lebih spesifik ya."
    db_add_message(chat_id, "assistant", fallback)
    return (fallback, media[:6])


async def maybe_summarize(chat_id):
    """When history grows large, fold older messages into a rolling summary and prune them."""
    if db_count(chat_id) < SUMMARIZE_AFTER:
        return
    history = db_recent(chat_id, 10_000)
    older = history[:-RECENT_MESSAGES] if len(history) > RECENT_MESSAGES else []
    if not older:
        return
    prev = db_get_summary(chat_id)
    convo = "\n".join(f"{m['role']}: {m['content']}" for m in older)
    prompt = (
        "Ringkas percakapan berikut jadi poin-poin penting yang perlu diingat untuk "
        "kelanjutan obrolan (fakta tentang user, preferensi, konteks, keputusan). "
        "Gabungkan dengan ringkasan lama bila ada. Tulis padat dalam Bahasa Indonesia.\n\n"
        f"Ringkasan lama:\n{prev or '(belum ada)'}\n\n"
        f"Percakapan yang perlu diringkas:\n{convo}"
    )
    try:
        data = await llm_chat([{"role": "user", "content": prompt}])
        msg = (data.get("choices") or [{}])[0].get("message", {}) or {}
        new_summary = (msg.get("content") or "").strip()
    except Exception as e:  # noqa: BLE001
        log.warning("summarize failed: %s", e)
        return
    if new_summary:
        db_set_summary(chat_id, new_summary)
        db_prune_older(chat_id, RECENT_MESSAGES)
        log.info("chat %s summarized; pruned to last %d messages", chat_id, RECENT_MESSAGES)


# ----------------------------- Telegram handlers -----------------------------
async def keep_typing(bot, chat_id, stop: asyncio.Event):
    """Send 'typing...' every few seconds until stopped (cloud replies can take a while)."""
    while not stop.is_set():
        try:
            await bot.send_chat_action(chat_id, constants.ChatAction.TYPING)
        except Exception:  # noqa: BLE001
            pass
        try:
            await asyncio.wait_for(stop.wait(), timeout=4)
        except asyncio.TimeoutError:
            pass


async def cmd_start(update: Update, context: ContextTypes.DEFAULT_TYPE):
    await update.message.reply_text(
        "Halo! Aku asisten AI kamu di Telegram 🤖\n\n"
        "Tanya apa aja — aku otomatis nyari info terkini di web kalau perlu, "
        "dan aku inget konteks obrolan kita.\n\n"
        "/reset — hapus ingatan obrolan ini",
        parse_mode=constants.ParseMode.MARKDOWN,
    )


async def cmd_reset(update: Update, context: ContextTypes.DEFAULT_TYPE):
    chat = update.effective_chat
    uid = update.effective_user.id if update.effective_user else 0
    # In groups the memory is shared, so only admins may wipe it.
    if chat.type in ("group", "supergroup") and not is_admin(uid):
        await update.message.reply_text("🔒 Di grup, /reset hanya bisa dilakukan admin.")
        return
    db_reset(chat.id)
    await update.message.reply_text("Oke, ingatan obrolan kita udah aku hapus. \U0001f9f9")


async def send_image(bot, chat_id, ref, caption=None):
    """Send an image (local path or http(s) URL) as a photo; for local files fall
    back to sending as a document. Web URLs that fail are skipped quietly."""
    is_url = isinstance(ref, str) and ref.startswith(("http://", "https://"))
    try:
        if is_url:
            await bot.send_photo(chat_id, photo=ref, caption=caption)
        else:
            with open(ref, "rb") as f:
                await bot.send_photo(chat_id, photo=f, caption=caption)
        return
    except Exception as e:  # noqa: BLE001
        if is_url:
            log.warning("failed web image %s: %s", ref, e)
            return
    try:
        with open(ref, "rb") as f:
            await bot.send_document(chat_id, document=f, caption=caption)
    except Exception as e:  # noqa: BLE001
        log.warning("failed to send image %s: %s", ref, e)


async def send_media_list(bot, chat_id, media):
    """Send a (ref, caption) list from generate_reply: images as photos, generated
    documents (txt/docx/pptx/pdf) as files."""
    for ref, caption in media:
        ext = os.path.splitext(ref)[1].lower() if isinstance(ref, str) else ""
        if ext in DOC_EXTS:
            try:
                with open(ref, "rb") as f:
                    await bot.send_document(chat_id, document=f, caption=caption)
            except Exception as e:  # noqa: BLE001
                log.warning("failed to send document %s: %s", ref, e)
        else:
            await send_image(bot, chat_id, ref, caption)


async def on_message(update: Update, context: ContextTypes.DEFAULT_TYPE):
    chat = update.effective_chat
    chat_id = chat.id
    msg = update.message
    raw = msg.text or ""
    bot_un = getattr(context.bot, "username", None)

    # In groups, only respond when the bot is addressed (mentioned or replied-to);
    # otherwise stay silent so we don't spam ordinary chatter.
    is_group = chat.type in ("group", "supergroup")
    mentioned = bool(bot_un and re.search(rf"@{re.escape(bot_un)}\b", raw, re.IGNORECASE))
    replied_to_bot = bool(
        msg.reply_to_message
        and msg.reply_to_message.from_user
        and msg.reply_to_message.from_user.id == context.bot.id
    )
    if is_group and not (mentioned or replied_to_bot):
        return

    # Strip the bot's @mention so the model sees a clean question.
    text = raw
    if bot_un:
        text = re.sub(rf"@{re.escape(bot_un)}\b", "", text, flags=re.IGNORECASE)
    text = text.strip()
    if not text:
        if raw.strip():  # addressed the bot with no actual question
            await msg.reply_text("Ya? Ada yang bisa aku bantu? 🙂")
        return

    stop = asyncio.Event()
    typing = asyncio.create_task(keep_typing(context.bot, chat_id, stop))
    images = []
    try:
        reply, images = await generate_reply(chat_id, text)
    except Exception as e:  # noqa: BLE001
        log.exception("error while generating reply")
        reply = f"Maaf, ada error: {e}"
    finally:
        stop.set()
        await typing

    # Telegram hard-limits messages to 4096 chars; split safely.
    for i in range(0, len(reply), 4000):
        await update.message.reply_text(reply[i : i + 4000])

    # Show any images/documents (web image-search, diagrams, generated docs).
    await send_media_list(context.bot, chat_id, images)


async def on_photo(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Photo sent normally (compressed, not as a file) -> vision: the model looks at
    it directly (base64) instead of just OCR'ing text out of it."""
    chat = update.effective_chat
    chat_id = chat.id
    msg = update.message
    bot_un = getattr(context.bot, "username", None)

    is_group = chat.type in ("group", "supergroup")
    caption = msg.caption or ""
    mentioned = bool(bot_un and re.search(rf"@{re.escape(bot_un)}\b", caption, re.IGNORECASE))
    replied_to_bot = bool(
        msg.reply_to_message
        and msg.reply_to_message.from_user
        and msg.reply_to_message.from_user.id == context.bot.id
    )
    if is_group and not (mentioned or replied_to_bot):
        return

    text = caption
    if bot_un:
        text = re.sub(rf"@{re.escape(bot_un)}\b", "", text, flags=re.IGNORECASE)
    text = text.strip() or "Jelasin/deskripsikan gambar ini."

    try:
        photo_file = await msg.photo[-1].get_file()
        raw = await photo_file.download_as_bytearray()
    except Exception as e:  # noqa: BLE001
        await msg.reply_text(f"⚠️ Gagal mengunduh gambar: {e}")
        return
    b64 = base64.b64encode(bytes(raw)).decode()

    # Save to disk too so edit_gambar can operate on it later in this chat.
    try:
        os.makedirs(STAGING_DIR, exist_ok=True)
        photo_path = os.path.join(STAGING_DIR, f"upload_{chat_id}_{int(time.time())}_{uuid.uuid4().hex[:6]}.jpg")
        with open(photo_path, "wb") as f:
            f.write(bytes(raw))
        register_last_image(chat_id, photo_path)
    except Exception as e:  # noqa: BLE001
        log.warning("failed to stage uploaded photo: %s", e)

    stop = asyncio.Event()
    typing = asyncio.create_task(keep_typing(context.bot, chat_id, stop))
    media = []
    try:
        reply, media = await generate_reply(chat_id, text, images=[b64])
    except Exception as e:  # noqa: BLE001
        log.exception("error while generating vision reply")
        reply = f"Maaf, ada error: {e}"
    finally:
        stop.set()
        await typing

    for i in range(0, len(reply), 4000):
        await update.message.reply_text(reply[i : i + 4000])
    await send_media_list(context.bot, chat_id, media)


MD = constants.ParseMode.MARKDOWN


async def cmd_ping(update: Update, context: ContextTypes.DEFAULT_TYPE):
    await update.message.reply_text("pong 🏓")


async def cmd_status(update: Update, context: ContextTypes.DEFAULT_TYPE):
    secs = int(time.time() - START_TIME)
    d, rem = divmod(secs, 86400)
    h, rem = divmod(rem, 3600)
    m, _ = divmod(rem, 60)
    uptime = f"{d}h {h}j {m}m" if d else (f"{h}j {m}m" if h else f"{m}m")
    await update.message.reply_text(
        f"🟢 *Status:* online\n"
        f"🔎 *Web search:* on\n"
        f"⏱️ *Uptime:* {uptime}",
        parse_mode=MD,
    )


async def cmd_name(update: Update, context: ContextTypes.DEFAULT_TYPE):
    await update.message.reply_text(INFO["name"], parse_mode=MD)


async def cmd_owner(update: Update, context: ContextTypes.DEFAULT_TYPE):
    await update.message.reply_text(INFO["owner"], parse_mode=MD)


async def cmd_skill(update: Update, context: ContextTypes.DEFAULT_TYPE):
    await update.message.reply_text(INFO["skill"])


async def cmd_myid(update: Update, context: ContextTypes.DEFAULT_TYPE):
    u = update.effective_user
    uid = u.id if u else 0
    role = "👑 admin" if is_admin(uid) else "👤 user"
    await update.message.reply_text(
        f"🆔 User ID kamu: `{uid}`\n"
        f"💬 Chat ID: `{update.effective_chat.id}`\n"
        f"🎚️ Peran: {role}",
        parse_mode=MD,
    )


async def cmd_whois(update: Update, context: ContextTypes.DEFAULT_TYPE):
    uid = update.effective_user.id if update.effective_user else 0
    if not is_admin(uid):
        await update.message.reply_text("🔒 Perintah ini khusus admin.")
        return
    target = update.message.reply_to_message
    if not target or not target.from_user:
        await update.message.reply_text(
            "ℹ️ Reply ke pesan orang yang mau dicek, lalu ketik /whois."
        )
        return
    u = target.from_user
    uname = f"@{u.username}" if u.username else "(tidak ada)"
    role = "👑 admin" if is_admin(u.id) else "👤 user"
    await update.message.reply_text(
        f"🆔 ID: `{u.id}`\n"
        f"👤 Nama: {u.full_name}\n"
        f"🔗 Username: {uname}\n"
        f"🎚️ Peran: {role}",
        parse_mode=MD,
    )




async def cmd_restart_bot(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Admin-only self-restart — no sudo/SSH needed. Sends itself SIGTERM (the same
    signal `systemctl restart` sends), which PTB handles as a graceful shutdown; since
    the systemd unit has Restart=always, the process comes back up on its own."""
    uid = update.effective_user.id if update.effective_user else 0
    if not is_admin(uid):
        await update.message.reply_text("🔒 Perintah ini khusus *admin*.", parse_mode=MD)
        return
    await update.message.reply_text(
        "🔄 Restarting bot... balik online dalam beberapa detik."
    )
    log.info("restart requested by admin %s via /restart_bot", uid)

    async def _delayed_restart():
        await asyncio.sleep(1)  # let the reply above actually flush to Telegram first
        os.kill(os.getpid(), signal.SIGTERM)

    asyncio.create_task(_delayed_restart())


async def ingest_chat_doc(chat_id, path, name, max_chunks=CHATDOC_MAX_INGEST):
    """Extract text from an uploaded file, embed its chunks, and store them scoped to
    this chat so the bot can answer about it. Returns (n_chunks_stored, n_chars)."""
    text = await asyncio.to_thread(doc_extract.extract, path)
    text = text or ""
    chunks = doc_extract.chunk_text(text)[:max_chunks]
    rows = []
    for idx, ch in enumerate(chunks):
        emb = await embed(ch, task="document")
        if emb is not None:
            rows.append((idx, ch, emb.tobytes()))
    db_add_chat_doc(chat_id, name, rows)
    return len(rows), len(text.strip())


async def on_document(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Anyone may upload a file -> the bot reads it (extract + embed) so it can be
    queried/summarized in this chat."""
    msg = update.message
    doc = msg.document if msg else None
    if doc is None:
        return
    chat_id = update.effective_chat.id
    fname = doc.file_name or "dokumen"
    ext = os.path.splitext(fname)[1].lower()
    if ext not in SUPPORTED_UPLOAD:
        await msg.reply_text(
            f"⚠️ Tipe file {ext or '?'} belum didukung.\n"
            "Didukung: PDF, Word, PowerPoint, Excel, txt, md, gambar (png/jpg)."
        )
        return
    if doc.file_size and doc.file_size > MAX_UPLOAD_MB * 1024 * 1024:
        await msg.reply_text(
            f"⚠️ File terlalu besar (~{doc.file_size // (1024 * 1024)} MB). "
            f"Batas upload via bot ~{MAX_UPLOAD_MB} MB."
        )
        return

    os.makedirs(STAGING_DIR, exist_ok=True)
    uid = update.effective_user.id if update.effective_user else 0
    staged = os.path.join(STAGING_DIR, f"{uid}_{int(time.time())}_{fname}")
    try:
        tg_file = await doc.get_file()
        await tg_file.download_to_drive(staged)
    except Exception as e:  # noqa: BLE001
        await msg.reply_text(f"⚠️ Gagal mengunduh file: {e}")
        return

    note = await msg.reply_text(f"📄 Lagi baca *{fname}*...", parse_mode=MD)
    try:
        n_chunks, _n_chars = await ingest_chat_doc(chat_id, staged, fname)
    except Exception as e:  # noqa: BLE001
        log.exception("chat-doc ingest failed")
        await note.edit_text(f"⚠️ Gagal membaca file: {e}")
        try:
            os.remove(staged)
        except OSError:
            pass
        return

    if n_chunks == 0:
        await note.edit_text(
            f"⚠️ Aku belum bisa menarik teks dari *{fname}* "
            "(mungkin gambar tanpa teks atau filenya kosong).",
            parse_mode=MD,
        )
        try:
            os.remove(staged)
        except OSError:
            pass
        return

    try:
        os.remove(staged)
    except OSError:
        pass
    await note.edit_text(
        f"📄 *{fname}* udah kebaca ({n_chunks} potongan).\n"
        "Mau dirangkum atau ada yang ditanya? Tinggal tanya aja di chat ini. 🙂",
        parse_mode=MD,
    )




# Command menus pushed to Telegram on startup. Default scope = everyone; each admin
# gets an extended menu in their own DM via a per-chat scope. (BotFather not needed.)
USER_COMMANDS = [
    BotCommand("start", "Mulai & info bot"),
    BotCommand("ping", "Cek bot aktif"),
    BotCommand("status", "Status bot"),
    BotCommand("myid", "Lihat ID & role kamu"),
    BotCommand("reset", "Hapus ingatan obrolan"),
    BotCommand("name", "Tentang bot"),
    BotCommand("owner", "Owner"),
    BotCommand("skill", "Kapabilitas bot"),
]
ADMIN_COMMANDS = USER_COMMANDS + [
    BotCommand("restart_bot", "[admin] Restart bot (tanpa SSH)"),
    BotCommand("whois", "[admin] Cek user ID (reply ke pesan)"),
    BotCommand("users", "[admin] User yang pernah chat bot"),
]


async def _set_command_menus(app):
    try:
        await app.bot.set_my_commands(USER_COMMANDS, scope=BotCommandScopeDefault())
        for aid in ADMIN_IDS:
            try:
                await app.bot.set_my_commands(ADMIN_COMMANDS, scope=BotCommandScopeChat(chat_id=aid))
            except Exception as e:  # noqa: BLE001
                log.warning("set admin menu for %s failed: %s", aid, e)
        log.info("command menus set (default + %d admin scope)", len(ADMIN_IDS))
    except Exception as e:  # noqa: BLE001
        log.warning("set_my_commands failed: %s", e)


async def _on_startup(app):
    # Drop any leftover staged files (older than a day) as a safety net.
    try:
        cutoff = time.time() - 86400
        if os.path.isdir(STAGING_DIR):
            for f in os.listdir(STAGING_DIR):
                p = os.path.join(STAGING_DIR, f)
                if os.path.isfile(p) and os.path.getmtime(p) < cutoff:
                    os.remove(p)
    except Exception as e:  # noqa: BLE001
        log.warning("staging cleanup failed: %s", e)
    await _set_command_menus(app)


async def track_user(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Runs for every update (own handler group) — records who interacted + when."""
    u = update.effective_user
    if u and not u.is_bot:
        try:
            db_touch_user(u.id, u.full_name, u.username)
        except Exception as e:  # noqa: BLE001
            log.warning("track_user failed: %s", e)


def _ago(ts):
    s = int(time.time() - ts)
    if s < 60:
        return "baru aja"
    if s < 3600:
        return f"{s // 60} mnt lalu"
    if s < 86400:
        return f"{s // 3600} jam lalu"
    return f"{s // 86400} hari lalu"


async def cmd_users(update: Update, context: ContextTypes.DEFAULT_TYPE):
    uid = update.effective_user.id if update.effective_user else 0
    if not is_admin(uid):
        await update.message.reply_text("🔒 Cuma admin.")
        return
    rows = db_list_users()
    if not rows:
        await update.message.reply_text("Belum ada user yang tercatat.")
        return
    lines = [f"👥 User yang pernah chat bot ({len(rows)}):"]
    for tid, name, uname, last, cnt in rows:
        who = ("@" + uname) if uname else (name or str(tid))
        lines.append(f"• {who} — {_ago(last)} ({cnt} pesan)")
    text = "\n".join(lines)
    for i in range(0, len(text), 4000):
        await update.message.reply_text(text[i:i + 4000])


def main():
    db_init()
    app = (
        Application.builder()
        .token(TELEGRAM_TOKEN)
        .concurrent_updates(True)
        .post_init(_on_startup)
        .build()
    )
    app.add_handler(CommandHandler("start", cmd_start))
    app.add_handler(CommandHandler("reset", cmd_reset))
    app.add_handler(CommandHandler("ping", cmd_ping))
    app.add_handler(CommandHandler("status", cmd_status))
    app.add_handler(CommandHandler("name", cmd_name))
    app.add_handler(CommandHandler("owner", cmd_owner))
    app.add_handler(CommandHandler("skill", cmd_skill))
    app.add_handler(CommandHandler("myid", cmd_myid))
    app.add_handler(CommandHandler("whois", cmd_whois))
    app.add_handler(CommandHandler("restart_bot", cmd_restart_bot))
    app.add_handler(CommandHandler("users", cmd_users))
    app.add_handler(TypeHandler(Update, track_user), group=-1)  # see every update, log who/when
    app.add_handler(MessageHandler(filters.Document.ALL, on_document))
    app.add_handler(MessageHandler(filters.PHOTO, on_photo))
    app.add_handler(MessageHandler(filters.TEXT & ~filters.COMMAND, on_message))
    log.info("Telegram Bot Kit starting (model=%s, polling)...", LLM_MODEL)
    app.run_polling(allowed_updates=Update.ALL_TYPES)


if __name__ == "__main__":
    main()
